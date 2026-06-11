"""通用媒体处理工具

提供文档文本提取、图片 base64 编码、多模态内容构建等功能。
不依赖特定平台，可被任何 Adapter 或 Pipeline 使用。

参考 nanobot 的设计模式：附件统一用本地文件路径表示，二进制数据不进入消息总线。
"""

from __future__ import annotations

import asyncio
import base64
import mimetypes
from abc import ABC, abstractmethod
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

# 可选依赖 — 顶层一次性 import，避免每次提取都 try/import
try:
    from pypdf import PdfReader as _PdfReader
except ImportError:
    _PdfReader = None

try:
    import docx as _docx_mod
except ImportError:
    _docx_mod = None  # type: ignore[assignment]

try:
    import openpyxl as _openpyxl_mod
except ImportError:
    _openpyxl_mod = None

try:
    from pptx import Presentation as _Presentation
except ImportError:
    _Presentation = None

# 图片扩展名集合
_IMAGE_EXTS = {
    ".png",
    ".jpg",
    ".jpeg",
    ".gif",
    ".bmp",
    ".webp",
    ".tif",
    ".tiff",
    ".ico",
    ".svg",
}


def _detect_image_mime(header: bytes) -> str | None:
    """通过 magic bytes 检测图片 MIME 类型。"""
    if header.startswith(b"\x89PNG\r\n\x1a\n"):
        return "image/png"
    if header.startswith(b"\xff\xd8\xff"):
        return "image/jpeg"
    if header.startswith(b"GIF87a") or header.startswith(b"GIF89a"):
        return "image/gif"
    if header.startswith(b"RIFF") and len(header) >= 12 and header[8:12] == b"WEBP":
        return "image/webp"
    if header.startswith(b"BM"):
        return "image/bmp"
    return None


def _is_image_file(path: str) -> bool:
    """判断文件是否为图片（通过扩展名和 magic bytes）。"""
    p = Path(path)
    if p.suffix.lower() in _IMAGE_EXTS:
        return True
    if not p.is_file():
        return False
    try:
        with open(p, "rb") as f:
            header = f.read(16)
        return _detect_image_mime(header) is not None
    except Exception:
        return False


# 纯文本文件扩展名集合
_TEXT_EXTS = {
    ".txt",
    ".md",
    ".csv",
    ".json",
    ".xml",
    ".html",
    ".htm",
    ".yaml",
    ".yml",
    ".toml",
    ".ini",
    ".cfg",
    ".py",
    ".js",
    ".ts",
    ".java",
    ".c",
    ".cpp",
    ".h",
    ".go",
    ".rs",
    ".rb",
    ".php",
    ".sh",
    ".bat",
    ".sql",
}


@dataclass(frozen=True)
class MediaBlock:
    """单个 LLM-ready 媒体内容块。

    由 MediaProcessor 根据文件类型生成，直接可嵌入 OpenAI 兼容的
    多模态 content 列表。不关心具体 LLM API 差异（由 LiteLLM 统一处理）。

    Attributes:
        block: OpenAI 兼容的 content block dict
        source_path: 原始本地文件路径（用于 _meta 元数据和调试）
        media_type: 媒体类型标签，如 "image"/"audio"/"video"/"document"
    """

    block: dict[str, Any]
    source_path: str
    media_type: str


@dataclass
class MediaProcessResult:
    """附件处理结果。

    Attributes:
        document_text: 从文档类文件提取的文本（已格式化）
        media_blocks: LLM-ready 媒体内容块列表（替代原 image_paths）
        att_meta: 每个附件的元数据列表
    """

    document_text: str = ""
    media_blocks: list[MediaBlock] = field(default_factory=list)
    att_meta: list[dict[str, Any]] = field(default_factory=list)


# ---------------------------------------------------------------------------
# MediaHandler 可插拔架构
# ---------------------------------------------------------------------------


class MediaHandler(ABC):
    """媒体类型处理器基类。

    每种媒体类型实现一个 Handler，由 MediaProcessor 统一调度。
    新增媒体类型只需注册新 Handler，无需修改核心流程。
    """

    @abstractmethod
    def can_handle(self, path: str) -> bool:
        """判断是否能处理该文件。

        应保持轻量（仅检查扩展名或 magic bytes），不做完整文件读取。
        """
        ...

    @abstractmethod
    async def encode(self, path: str) -> MediaBlock | None:
        """将文件编码为 LLM-ready content block。

        涉及文件 I/O 时应使用 asyncio.to_thread() 避免阻塞事件循环。
        """
        ...


class ImageHandler(MediaHandler):
    """图片处理器 → image_url block"""

    _EXTS = _IMAGE_EXTS

    def can_handle(self, path: str) -> bool:
        return Path(path).suffix.lower() in self._EXTS or _is_image_file(path)

    async def encode(self, path: str) -> MediaBlock | None:
        p = Path(path)
        if not p.is_file():
            return None
        try:
            raw = await asyncio.to_thread(p.read_bytes)
            mime = _detect_image_mime(raw) or mimetypes.guess_type(path)[0]
            if not mime or not mime.startswith("image/"):
                return None
            b64 = base64.b64encode(raw).decode()
            return MediaBlock(
                block={
                    "type": "image_url",
                    "image_url": {"url": f"data:{mime};base64,{b64}"},
                    "_meta": {"path": str(p)},
                },
                source_path=str(p),
                media_type="image",
            )
        except Exception:
            return None


class MediaProcessor:
    """媒体处理器 - 统一处理附件（分离图片/文档、提取文本、编码图片）。

    参考 nanobot ContextBuilder 的设计，将附件处理封装为可复用的类。

    Example:
        processor = MediaProcessor()
        result = await processor.process(["/path/to/doc.pdf", "/path/to/image.png"])
        content = processor.build_content("用户问题", result.media_blocks)
    """

    def __init__(self, max_file_size: int = 50 * 1024 * 1024) -> None:
        self.max_file_size = max_file_size
        self._handlers: list[MediaHandler] = [
            ImageHandler(),
        ]

    async def _try_encode_as_media(self, path: str) -> tuple[MediaBlock | None, bool]:
        """尝试按优先级匹配媒体类型并编码。

        Returns:
            (MediaBlock | None, matched): 编码结果 + 是否有 handler 匹配
        """
        for handler in self._handlers:
            if handler.can_handle(path):
                block = await handler.encode(path)
                return block, True
        return None, False

    async def process(self, attachments: list[str]) -> MediaProcessResult:
        """处理附件列表，分离媒体和文档，提取文档文本。

        Args:
            attachments: 本地文件路径列表

        Returns:
            MediaProcessResult: 处理结果，包含文档文本、媒体块和元数据
        """
        doc_texts: list[str] = []
        media_blocks: list[MediaBlock] = []
        att_meta: list[dict[str, Any]] = []

        for path in attachments:
            p = Path(path)
            meta = {"path": path, "exists": p.is_file()}

            if not p.is_file():
                meta["error"] = "file not found"
                att_meta.append(meta)
                continue

            # 尝试按优先级匹配媒体类型
            block, matched = await self._try_encode_as_media(path)
            if block is not None:
                media_blocks.append(block)
                meta["type"] = block.media_type
                att_meta.append(meta)
                continue

            # handler 匹配但编码失败（文件损坏/不可读），跳过文档提取
            if matched:
                meta["type"] = "media_error"
                meta["error"] = f"Media handler failed for {Path(path).name}"
                att_meta.append(meta)
                continue

            # 非媒体文件：提取文档文本（在线程池中执行避免阻塞事件循环）
            meta["type"] = "document"
            extracted = await asyncio.to_thread(self.extract_document_text, path)
            if extracted and not extracted.startswith("["):
                doc_texts.append(f"[File: {p.name}]\n{extracted}")
                meta["extracted"] = True
            elif extracted:
                doc_texts.append(extracted)
                meta["extracted"] = False
                meta["fallback"] = extracted
            else:
                meta["extracted"] = False

            att_meta.append(meta)

        document_text = "\n\n".join(doc_texts) if doc_texts else ""
        return MediaProcessResult(
            document_text=document_text,
            media_blocks=media_blocks,
            att_meta=att_meta,
        )

    def build_content(
        self,
        text_content: str,
        media_blocks: list[MediaBlock],
    ) -> str | list[dict[str, Any]]:
        """构建用户消息的 content 字段。

        如果没有媒体块，返回纯文本字符串。
        如果有媒体块，返回 OpenAI 兼容的多模态列表（image_url blocks + text block）。

        Args:
            text_content: 文本内容
            media_blocks: MediaBlock 列表

        Returns:
            str | list[dict]: 纯文本或 OpenAI 兼容的多模态 content
        """
        blocks = [mb.block for mb in media_blocks]
        if not blocks:
            return text_content
        return blocks + [{"type": "text", "text": text_content}]

    # ------------------------------------------------------------------
    # Document text extraction (synchronous, called via to_thread)
    # ------------------------------------------------------------------

    def extract_document_text(self, path: str) -> str | None:
        """从文档中提取文本内容。

        支持格式：PDF, DOCX, XLSX, PPTX, TXT, MD, CSV, JSON, XML, HTML, YAML, TOML, INI
        可选依赖未安装时返回提示文本。
        """
        p = Path(path)
        if not p.is_file():
            return None

        if p.stat().st_size > self.max_file_size:
            return f"[File too large: {p.name} ({p.stat().st_size / 1024 / 1024:.1f}MB)]"

        ext = p.suffix.lower()

        if ext in _TEXT_EXTS:
            return self._read_text_file(p)

        if ext == ".pdf":
            return self._extract_pdf(p)

        if ext == ".docx":
            return self._extract_docx(p)

        if ext == ".xlsx":
            return self._extract_xlsx(p)

        if ext == ".pptx":
            return self._extract_pptx(p)

        return f"[File: {p.name} (unsupported format: {ext})]"

    @staticmethod
    def _read_text_file(p: Path) -> str:
        try:
            return p.read_text(encoding="utf-8")
        except UnicodeDecodeError:
            try:
                return p.read_text(encoding="gbk")
            except Exception:
                return f"[Text file: {p.name} (encoding error)]"

    @staticmethod
    def _extract_pdf(p: Path) -> str:
        if _PdfReader is None:
            return f"[PDF file: {p.name} (install pypdf to extract text)]"
        try:
            reader = _PdfReader(str(p))
            texts = []
            for page in reader.pages:
                text = page.extract_text()
                if text:
                    texts.append(text)
            return "\n".join(texts) if texts else f"[PDF: {p.name} (no extractable text)]"
        except Exception as e:
            return f"[PDF: {p.name} (extraction error: {e})]"

    @staticmethod
    def _extract_docx(p: Path) -> str:
        if _docx_mod is None:
            return f"[DOCX file: {p.name} (install python-docx to extract text)]"
        try:
            document = _docx_mod.Document(str(p))
            texts = [para.text for para in document.paragraphs if para.text.strip()]
            return "\n".join(texts) if texts else f"[DOCX: {p.name} (empty)]"
        except Exception as e:
            return f"[DOCX: {p.name} (extraction error: {e})]"

    @staticmethod
    def _extract_xlsx(p: Path) -> str:
        if _openpyxl_mod is None:
            return f"[XLSX file: {p.name} (install openpyxl to extract text)]"
        try:
            wb = _openpyxl_mod.load_workbook(str(p), data_only=True)
            texts = []
            for sheet in wb.worksheets:
                sheet_texts = []
                for row in sheet.iter_rows(values_only=True):
                    row_text = "\t".join(str(cell) for cell in row if cell is not None)
                    if row_text.strip():
                        sheet_texts.append(row_text)
                if sheet_texts:
                    texts.append(f"--- Sheet: {sheet.title} ---\n" + "\n".join(sheet_texts))
            return "\n\n".join(texts) if texts else f"[XLSX: {p.name} (empty)]"
        except Exception as e:
            return f"[XLSX: {p.name} (extraction error: {e})]"

    @staticmethod
    def _extract_pptx(p: Path) -> str:
        if _Presentation is None:
            return f"[PPTX file: {p.name} (install python-pptx to extract text)]"
        try:
            prs = _Presentation(str(p))
            texts = []
            for i, slide in enumerate(prs.slides, 1):
                slide_texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_texts.append(shape.text.strip())
                if slide_texts:
                    texts.append(f"--- Slide {i} ---\n" + "\n".join(slide_texts))
            return "\n\n".join(texts) if texts else f"[PPTX: {p.name} (empty)]"
        except Exception as e:
            return f"[PPTX: {p.name} (extraction error: {e})]"
